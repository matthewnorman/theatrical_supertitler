import pptx
import argparse

SKIP_STARTS = (
    'PROLOGUE', 'ACT', 'Scene One', 'Scene Two',
    'Place', 'English Translation',
    'END', 'A country cross', 'etc.',
    'Intermezzo', 'Same scene as before',
    'THE PLAY',
)
LINES_PER_SLIDE = 3


def add_slide(slide_lines, prs, slide_layout):
    if slide_lines == []:
        # No need to put any text here
        # because there's been nothing since
        # the last line break.
        return

    slide = prs.slides.add_slide(slide_layout)
    top = height = pptx.util.Inches(1)
    left = pptx.util.Inches(1)
    width = pptx.util.Inches(8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()

    p.text = ''.join(slide_lines)
    p.font.size = pptx.util.Pt(20)
    p.font.color.rgb = pptx.dml.color.RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = pptx.enum.text.PP_ALIGN.CENTER


def run(input_script, output_file):
    with open(input_script, 'r') as input:
        input_lines = input.readlines()

    prs = pptx.Presentation()

    slide_layout = prs.slide_layouts[6]

    slide_lines = []
    parenthesis_flag = False
    for line in input_lines:
        if line == '\n' or line.startswith(SKIP_STARTS):
            add_slide(slide_lines=slide_lines,
                      prs=prs,
                      slide_layout=slide_layout)
            slide_lines = []
            continue

        if line.startswith('('):
            parenthesis_flag = True

        if parenthesis_flag and not ')' in line:
            # We're still in the comments.
            continue

        if parenthesis_flag and ')' in line:
            # Exit the comments
            parenthesis_flag = False
            continue

        # If we're here, that means we are adding
        # onto a continuing speech.
        slide_lines.append(line)

        if len(slide_lines) >= LINES_PER_SLIDE:
            add_slide(slide_lines=slide_lines,
                      prs=prs,
                      slide_layout=slide_layout)
            print(slide_lines)
            slide_lines = []

    add_slide(slide_lines=slide_lines,
              prs=prs,
              slide_layout=slide_layout)
    slide_lines = []

    prs.save(output_file)
    


if __name__ == '__main__':

    parser = argparse.ArgumentParser(
        description='Process text script into pptx supertitles'
    )

    parser.add_argument('input_script', type=str,
                        help='Input text file with the script')
    parser.add_argument('output_file', type=str,
                        default='final.pptx',
                        help='Output powerpoint destination')

    args = parser.parse_args()

    run(
        input_script=args.input_script,
        output_file=args.output_file
    )
